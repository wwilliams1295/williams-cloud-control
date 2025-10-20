# cloud.py — Williams GPT-5 Local Node with GV SMS via Gmail bridge
# - Python 3.7+ compatible (uses Optional[...] instead of |)
# - GV-safe output (no truncation on blank lines; markdown stripped)
# - Robust allowlists (E.164 + last-10 matching); includes Bill (+1 561-371-6077)
# - Ignores Google Voice admin notices (GV_IGNORE_ADMIN), and settings-link hard-ignore (GV_HARD_IGNORE_SETTINGS)
# - Gmail API send with SMTP fallback and final file-write safety
# - Strips GV boilerplate from inbound messages so only the SMS text is processed
# - Extra debug logs when DEBUG_LOG=1
# - ALWAYS replies to simple greetings ("hello"/"hi"/"hey"/"ping")
# - Optional GREET_FORCE=1 shows greeting on every message while testing
# - Never sends empty messages; falls back to "ping" payload if cleaned body is empty

from fastapi import FastAPI, Request, HTTPException, Form, Body
from fastapi.responses import PlainTextResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from typing import Optional, List, Tuple, Dict, Any, Callable, Deque
import os
import sys
import json
import hmac
import hashlib
import requests
import random
import secrets
import string
import traceback
import time
import mimetypes
import base64
import asyncio
import threading
import re
import subprocess
from datetime import datetime
from zoneinfo import ZoneInfo
from dotenv import load_dotenv
from bs4 import BeautifulSoup
from duckduckgo_search import DDGS
from pathlib import Path
from urllib.parse import quote as urlquote, urljoin

# Excel / Files / PPT
from openpyxl import Workbook
from pptx import Presentation

# Email (SMTP replies)
import smtplib
from email.message import EmailMessage
from email.utils import parseaddr
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders

# LLM router (your existing)
from agent import superchat

# Optional (SEC / data)
import pandas as pd
from collections import deque

# =============== Google Voice via Gmail (OAuth user data) ===============
try:
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow  # noqa: F401 (not used directly)
    from googleapiclient.discovery import build
    from google.auth.transport.requests import Request as GoogleRequest

    _GOOGLE_OK = True
except Exception:
    _GOOGLE_OK = False

# ==================== APP INIT ====================
load_dotenv()
app = FastAPI()

# ==================== CONFIG ======================
SECRET: str = os.getenv("DISPATCH_SECRET", "dev-secret")
MAC_AGENT_URL: str = os.getenv("MAC_AGENT_URL", "http://127.0.0.1:8787")
WIN_AGENT_URL: Optional[str] = os.getenv("WIN_AGENT_URL")  # optional
TZ: str = os.getenv("APP_TIMEZONE", "America/Chicago")

DB_PATH: str = os.getenv("CODES_DB", os.path.expanduser("~/jarvis-demo/commands.json"))
SHEETS_DIR: str = os.getenv(
    "SHEETS_DIR", os.path.expanduser("~/Documents/JarvisSheets")
)
Path(SHEETS_DIR).mkdir(parents=True, exist_ok=True)

# Cloud output dir for generated files
CLOUD_OUT_DIR = Path(
    os.getenv("CLOUD_OUT_DIR", os.path.expanduser("~/Documents/JarvisCloud"))
)
CLOUD_OUT_DIR.mkdir(parents=True, exist_ok=True)

# Serve files publicly (for SMS links)
PUBLIC_BASE_URL: str = os.getenv("PUBLIC_BASE_URL", "http://127.0.0.1:8000").rstrip("/")
app.mount("/files", StaticFiles(directory=str(CLOUD_OUT_DIR)), name="files")

# Email config (SMTP path)
EMAIL_FROM: str = os.getenv("EMAIL_FROM", "")
EMAIL_WHITELIST: List[str] = [
    e.strip().lower()
    for e in (os.getenv("EMAIL_WHITELIST", "").split(","))
    if e.strip()
]
SMTP_HOST: str = os.getenv("SMTP_HOST", "")
SMTP_PORT: int = int(os.getenv("SMTP_PORT", "587"))
SMTP_USER: str = os.getenv("SMTP_USER", "")
SMTP_PASS: str = os.getenv("SMTP_PASS", "")
SMTP_DEBUG: bool = os.getenv("SMTP_DEBUG", "0") == "1"
COMPANY_ADDRESS: str = os.getenv("COMPANY_ADDRESS", "")

# Greeting
GREET_COOLDOWN_SECONDS: int = int(os.getenv("GREET_COOLDOWN_SECONDS", "300"))
_last_greet_at: Dict[str, float] = {}

# Node name
NODE_NAME: str = os.getenv("NODE_NAME", "Echo-Nine")

# === Gmail / GV files & scopes ===
DATA_DIR = Path(os.getenv("DATA_DIR", "."))
GV_CLIENT_JSON = str(DATA_DIR / "client_secret.json")
GV_TOKEN_JSON = str(DATA_DIR / "token.json")
GMAIL_SCOPES = [
    "https://www.googleapis.com/auth/gmail.modify",
    "https://www.googleapis.com/auth/gmail.send",
]
# Prefer real message notifications; still accept @txt.voice
GMAIL_QUERY: str = os.getenv(
    "GMAIL_QUERY",
    'in:inbox is:unread newer_than:7d (from:@txt.voice.google.com OR (from:voice-noreply@google.com subject:"New message"))',
)
POLL_SECONDS: int = int(os.getenv("POLL_SECONDS", "20"))

# ==================== PLUGINS (auto-capabilities) ==========
# 3.7-safe typing for conditional imports
_get_plugin: Callable[[str], Any]
_list_plugins: Callable[[], List[str]]
try:
    from plugins.loader import (
        load_all as _load_plugins,
        get as _get_plugin,
        list_plugins as _list_plugins,
    )

    _PLUGIN_REG = _load_plugins()
except Exception as _e:
    _PLUGIN_REG = {}

    def _get_plugin(name: str) -> Any:
        raise KeyError(f"plugins not available: {name} ({_e})")

    def _list_plugins() -> List[str]:
        return []


# ==================== UTIL ========================
PHONE_RE = re.compile(r"(?:\+?1)?\D?(\d{3})\D?(\d{3})\D?(\d{4})")


def normalize_phone(raw: str) -> Optional[str]:
    if not raw:
        return None
    s = re.sub(r"\D+", "", raw)
    if s.startswith("1") and len(s) == 11:
        return f"+{s}"
    if len(s) == 10:
        return f"+1{s}"
    if (raw or "").startswith("+") and len(s) >= 10:
        return f"+{s}"
    m = PHONE_RE.search(raw or "")
    if m:
        return f"+1{m.group(1)}{m.group(2)}{m.group(3)}"
    return None


def last10_digits(s: str) -> Optional[str]:
    if not s:
        return None
    d = re.sub(r"\D+", "", s)
    return d[-10:] if len(d) >= 10 else None


def _is_gv_gateway(addr: str) -> bool:
    a = (addr or "").lower()
    return ("@txt.voice.google.com" in a) or ("voice-noreply@google.com" in a)


def _is_email_allowed(
    addr: str, whitelist: List[str], also: Dict[str, dict], allow_gv: bool = True
) -> bool:
    addr = (addr or "").strip().lower()
    if not addr:
        return False
    wl = [w.strip().lower() for w in (whitelist or []) if w.strip()]
    if not wl:
        return True
    if addr in wl:
        return True
    for w in wl:
        if w.startswith("*@"):
            dom = w[2:]
            if addr.endswith("@" + dom):
                return True
    if addr in (k.strip().lower() for k in also.keys()):
        return True
    if allow_gv and _is_gv_gateway(addr):
        return True
    return False


def extract_phone_from_gv_headers(
    headers: Optional[Dict[str, str]], subject: str, body: str
) -> Optional[str]:
    """
    GV local-part like: 19177194526.15613891295.CRmvhhlGR7
    Extract sender's phone number from subject line first, then from email headers.
    """
    headers = headers or {}
    subject = subject or ""
    body = body or ""

    # First, try to extract phone number from subject line
    # Subject format: "New text message from (561) 389-1295"
    subject_match = re.search(r'\((\d{3})\)\s*(\d{3})-(\d{4})', subject)
    if subject_match:
        area_code, prefix, number = subject_match.groups()
        phone_str = f"+1{area_code}{prefix}{number}"
        normalized = normalize_phone(phone_str)
        if normalized:
            return normalized

    def pick_from_local(local: str) -> Optional[str]:
        if not local:
            return None
        tokens = [t for t in local.split(".") if t]
        # First, look for sender numbers (not the GV number)
        for tok in tokens:
            n = normalize_phone(tok)
            if not n:
                continue
            f10 = last10_digits(n)
            # Skip if this is the GV number itself
            if f10 == "9177194526":
                continue
            if f10 and f10 in ALLOWED_NUMBERS_LAST10:
                return n
        for tok in tokens:
            n = normalize_phone(tok)
            if n:
                return n
        return None

    for key in ("reply-to", "from"):
        v = (headers.get(key) or "").strip()
        if v:
            m = re.search(r"<([^>]+)>", v)
            addr = m.group(1) if m else v
            local = addr.split("@", 1)[0]
            cand = pick_from_local(local)
            if cand:
                return cand

    for key in ("reply-to", "from", "delivered-to", "to"):
        v = headers.get(key, "")
        n = normalize_phone(v)
        if n:
            return n

    for blob in ((subject or ""), (body or "")):
        n = normalize_phone(blob)
        if n:
            return n
    return None


def now_str() -> str:
    try:
        return datetime.now(ZoneInfo(TZ)).strftime("%A, %B %-d, %Y %I:%M %p %Z")
    except Exception:
        return datetime.now().strftime("%A, %B %d, %Y %I:%M %p")


def twiml(msg: str) -> PlainTextResponse:
    return PlainTextResponse(
        f'<?xml version="1.0"?><Response><Message>{msg}</Message></Response>',
        media_type="application/xml",
    )


def sign_payload(d: dict) -> str:
    return hmac.new(
        SECRET.encode(), json.dumps(d, sort_keys=True).encode(), hashlib.sha256
    ).hexdigest()


def maybe_greeting(sender_key: str, name: str) -> str:
    # Force greeting on every message while testing if GREET_FORCE=1
    if os.getenv("GREET_FORCE", "0") == "1":
        call_signs = [
            "Williams Echo-Nine",
            "Williams Core-One",
            "Palm Node",
            "Houston Command",
        ]
        rng = random.SystemRandom()
        call_sign = rng.choice(call_signs)
        return f"{name} — You have been authenticated to the Williams Secured Cloud Control Server. Uplink established with {call_sign}. How can I assist you today?"
    now = time.time()
    last = _last_greet_at.get(sender_key, 0.0)
    if now - last < GREET_COOLDOWN_SECONDS:
        return ""
    _last_greet_at[sender_key] = now
    call_signs = [
        "Williams Echo-Nine",
        "Williams Core-One",
        "Palm Node",
        "Houston Command",
    ]
    rng = random.SystemRandom()
    call_sign = rng.choice(call_signs)
    # INLINE to avoid GV truncation at blank lines
    return (
        f"{name} — You have been authenticated to the Williams Secured Cloud Control Server. "
        f"Uplink established with {call_sign}. How can I assist you today?"
    )


def jarvis_signature() -> str:
    return (
        "\n\nBest regards,\n"
        f"Jarvis ({NODE_NAME})\n"
        "Williams Secured Cloud Control Server AI"
    )


def _public_file_url(path: str) -> str:
    return f"{PUBLIC_BASE_URL}/files/{urlquote(os.path.basename(path))}"


def is_gv_admin_notice(
    subject: str, body: str, headers: Optional[Dict[str, str]] = None
) -> bool:
    # Heuristic admin filter is gated by env; set GV_IGNORE_ADMIN=0 to bypass during testing
    if os.getenv("GV_IGNORE_ADMIN", "1") != "1":
        return False

    subj = (subject or "").lower()
    bod = (body or "").lower()
    hdr_from = ((headers or {}).get("from") or "").lower()
    hdr_rt = ((headers or {}).get("reply-to") or "").lower()
    blob = " ".join([subj, bod, hdr_from, hdr_rt])

    # Only treat as admin if clearly from GV system and matches admin phrases
    is_from_gv = ("voice-noreply@google.com" in blob) or (
        "@txt.voice.google.com" in blob
    )

    admin_markers = [
        "notification settings",
        "messaging settings",
        "/settings#messaging",
        "settings have been updated",
        "message notifications are off",
        "you will no longer receive",
        "turn on notifications",
        "policy update",
        "service update",
        "google voice terms",
        "security alert",
    ]
    sms_markers = ["new message", "sent a message", "text message"]

    if not is_from_gv:
        return False
    if any(p in blob for p in sms_markers):
        return False
    return any(p in blob for p in admin_markers)


def make_gv_friendly(text: str) -> str:
    """Strip Markdown, collapse blanks, normalize bullets, soft-cap for GV."""
    if not text:
        return ""
    s = text
    s = re.sub(r"\*\*(.*?)\*\*", r"\1", s)
    s = re.sub(r"__(.*?)__", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    s = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", s)
    s = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1: \2", s)
    s = re.sub(r"^#{1,6}\s*", "", s, flags=re.M)
    s = re.sub(
        r"\n{2,}Best regards,.*?\Z", "", s, flags=re.S | re.I
    )  # drop signature for SMS
    s = re.sub(r"\n{2,}", "\n", s)
    s = s.replace("• ", "- ").replace("▸ ", "- ")
    s = s.strip()
    MAX = int(os.getenv("GV_SOFT_CHAR_LIMIT", "1000"))
    if len(s) > MAX:
        s = s[: MAX - 3].rstrip() + "…"
    return s


def needs_freshness(q: str) -> bool:
    q = (q or "").lower()
    trigger_words = [
        "trend",
        "trends",
        "today",
        "now",
        "latest",
        "this year",
        "this month",
        "update",
        "news",
        "forecast",
        "outlook",
    ]
    return any(t in q for t in trigger_words)


# ----- Strip GV boilerplate from inbound messages
URL_LINE = re.compile(r"^\s*<?https?://", re.I)
NOISE_PHRASES = [
    "your account",
    "help center",
    "google voice",
    "notification settings",
    "support.google.com/voice",
    "voice.google.com",
    "productforums.google.com",
    "txt.voice.google.com",
    "from: google",
    "to:",
    "from:",
    "date:",
]

# NEW: subject stripper + admin blob checker
SUBJECT_LINE = re.compile(r"^\s*subject\s*:\s*", re.I)
ADMIN_LIKE_PHRASES = [
    "notification preferences",
    "notification settings",
    "email notifications for text messages",
    "unsubscribe from these notifications",
    "update your notification settings",
    "this message serves as a reminder",
    "your satisfaction is paramount",
]


def strip_subject_lines(s: str) -> str:
    if not s:
        return ""
    lines = [l for l in s.splitlines() if not SUBJECT_LINE.match(l)]
    return "\n".join(lines).strip()


def looks_like_admin_blob(s: str) -> bool:
    """Heuristic: treat GV admin/notification blobs as non-SMS."""
    if not s:
        return False
    low = s.lower()
    admin_markers = [
        "notification preferences",
        "messaging settings",
        "update notification settings",
        "/settings#messaging",
        "you will no longer receive",
        "turn on notifications",
        "policy update",
        "service update",
        "google voice terms",
        "security alert",
    ]
    # Only classify as admin if the text is short or heavily matches admin phrases
    if len(low) < 100 and sum(p in low for p in admin_markers) >= 2:
        return True
    return False


def extract_gv_sms_payload(raw: str) -> str:
    """Remove GV boilerplate and keep the actual SMS text."""
    if not raw:
        return ""
    lines = [l.strip().strip("<>").strip() for l in raw.splitlines() if l.strip()]
    clean: List[str] = []
    for l in lines:
        low = l.lower()
        if URL_LINE.match(l):
            continue
        if any(p in low for p in NOISE_PHRASES):
            continue
        if re.search(r"[A-Za-z0-9]", l):
            clean.append(l)
    if not clean:
        s = raw.strip()
    else:
        s = "\n".join(clean[:3]).strip()
    # NEW: kill any Subject: lines that slip into bodies
    s = strip_subject_lines(s)
    return s


# ==================== ALLOW LISTS ====================
# SMS allowed numbers
_raw_allowed_numbers: Dict[str, Dict[str, str]] = {}
# From env (comma-separated): +15613891295,+11234567890
_env_nums = [
    n.strip() for n in os.getenv("ALLOWED_NUMBERS", "").split(",") if n.strip()
]
for n in _env_nums:
    _raw_allowed_numbers[n] = {"name": "Guest", "role": "guest"}

# Hardcode authorized users
_raw_allowed_numbers["+15613891295"] = {"name": "Chris Williams", "role": "admin"}
_raw_allowed_numbers["+15613716077"] = {"name": "Bill", "role": "admin"}  # 561-371-6077
_raw_allowed_numbers["+19177194526"] = {"name": "Google Voice User", "role": "admin"}  # 917-719-4526
_raw_allowed_numbers["+6822055698"] = {
    "name": "Angry Munch",
    "role": "admin",
}  # example extra

ALLOWED_NUMBERS: Dict[str, Dict[str, str]] = {}
for k, v in _raw_allowed_numbers.items():
    nk = normalize_phone(k)
    if nk:
        ALLOWED_NUMBERS[nk] = v

ALLOWED_NUMBERS_LAST10 = {
    last10_digits(nk) for nk in ALLOWED_NUMBERS.keys() if last10_digits(nk)
}

if os.getenv("DEBUG_LOG", "0") == "1":
    print("[BOOT] ALLOWED_NUMBERS =", ALLOWED_NUMBERS)
    print("[BOOT] ALLOWED_NUMBERS_LAST10 =", ALLOWED_NUMBERS_LAST10)

# Email allow list (explicit map)
ALLOWED_EMAILS: Dict[str, Dict[str, str]] = {
    "william.c.williams@outlook.com": {"name": "Chris Williams", "role": "admin"},
}


# ==================== EMAIL / FILE HELPERS ==========
def send_email(
    to_addr: str, subject: str, body: str, attachments: Optional[List[str]] = None
) -> str:
    if not (SMTP_HOST and SMTP_PORT and SMTP_USER and SMTP_PASS and EMAIL_FROM):
        return "(Email disabled: SMTP settings missing)"
    if COMPANY_ADDRESS:
        body = f"{body}\n\n—\n{COMPANY_ADDRESS}"

    to_email = parseaddr(to_addr)[1] or to_addr
    msg = EmailMessage()
    msg["From"] = EMAIL_FROM
    msg["To"] = to_email
    msg["Subject"] = subject
    msg.set_content(body)

    for path in attachments or []:
        if not os.path.isfile(path):
            continue
        ctype, _ = mimetypes.guess_type(path)
        maintype, subtype = (ctype or "application/octet-stream").split("/", 1)
        with open(path, "rb") as f:
            msg.add_attachment(
                f.read(),
                maintype=maintype,
                subtype=subtype,
                filename=os.path.basename(path),
            )

    try:
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=20) as s:
            if SMTP_DEBUG:
                s.set_debuglevel(1)
            s.ehlo()
            s.starttls()
            s.ehlo()
            s.login(SMTP_USER, SMTP_PASS)
            s.send_message(msg, from_addr=EMAIL_FROM, to_addrs=[to_email])
        return "Email sent."
    except Exception as e:
        return f"Email send error: {e}"


def _extract_paths(text: str) -> List[str]:
    paths: List[str] = []
    for line in (text or "").splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith(("PPTX:", "PDF:", "XLSX:")):
            p = line.split(":", 1)[1].strip()
            if os.path.isfile(p):
                paths.append(p)
    return paths


def _replace_paths_with_links(reply_text: str, paths: List[str]) -> str:
    out = reply_text or ""
    for p in paths:
        url = _public_file_url(p)
        out = re.sub(
            rf"(?mi)^\s*(PDF|PPTX|XLSX)\s*:\s*{re.escape(p)}\s*$", f"Link: {url}", out
        )
        if url not in out:
            out += f"\nLink: {url}"
    return out


# ==================== LIGHT MEMORY =================
MEMORY: Deque[Dict[str, str]] = deque(maxlen=20)


def remember(user_text: str, ai_text: str) -> None:
    MEMORY.append({"user": user_text, "ai": ai_text})


def memory_text() -> str:
    if not MEMORY:
        return "(no prior context)"
    lines: List[str] = []
    for m in MEMORY:
        lines.append(f"User: {m['user']}")
        lines.append(f"AI: {m['ai']}")
    return "\n".join(lines[-12:])


# ==================== LLM HELPERS ==================
async def ask_llm_async(prompt: str, name: str = "Operator") -> str:
    # NOTE: steer SMS away from email formatting
    system = (
        "You are Jarvis, the Williams Secured Cloud Control Server AI (codename Echo-Nine). "
        "For SMS, reply in 1–4 short lines, conversational and direct, but also detailed and sophisticated. "
        "Do NOT include email headers like Subject/To/From or long salutations. "
        "For documents/decks, you may switch to a formal executive style with premium bullets."
    )
    formatted_prompt = f"{prompt.strip()}\n\n(Operator: {name})"
    return await superchat(formatted_prompt, system=system)


def ask_llm(prompt: str, name: str = "Operator") -> str:
    coro = ask_llm_async(prompt, name)
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        return asyncio.run(coro)

    result: Dict[str, Any] = {"val": None, "err": None}

    def _runner() -> None:
        try:
            result["val"] = asyncio.run(coro)
        except Exception as e:
            result["err"] = e

    t = threading.Thread(target=_runner, daemon=True)
    t.start()
    t.join()
    if result["err"] is not None:
        return f"(Router error) {result['err']}"
    return str(result["val"] or "")


# ==================== WEB / RESEARCH ================
def ddg_search(query: str, max_results: int = 5) -> List[Dict[str, Any]]:
    try:
        ddgs = DDGS()
        return list(ddgs.text(query, max_results=max_results))
    except Exception as e:
        return [{"title": "Web search error", "body": str(e), "href": ""}]


def format_search_results(results: List[Dict[str, Any]]) -> str:
    if not results:
        return ""
    out: List[str] = []
    for i, r in enumerate(results, start=1):
        out.append(
            f"({i}) {r.get('title', 'Untitled')}\n{r.get('body', '')}\n🔗 {r.get('href', '')}"
        )
    return "\n\n".join(out)


def fetch_article(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Williams Cloud Node)"}
        r = requests.get(url, headers=headers, timeout=12)
        soup = BeautifulSoup(r.text, "html.parser")
        text = " ".join(p.get_text(" ", strip=True) for p in soup.find_all("p"))[:6000]
        if not text.strip():
            return "The page content appears empty or blocked. Try another URL."
        return ask_llm(f"Summarize clearly in 4–6 premium bullets:\n{text}")
    except Exception as e:
        return f"Article fetch error: {e}"


def research_with_perplexity(query: str, name: str) -> str:
    prompt = (
        "use: perplexity\n"
        "Research the web and produce a concise brief in 5–7 premium bullets (• / ▸). "
        "Do not include subject headers, footnote numbers or inline source names with URLs at the end of relevant bullets."
        "Use incredibly detailed numbers like an investment banker, be professional, be quantative, and use less than or equal to 1,000 characters."
        f"\n\nQuery: {query}"
    )
    ans = ask_llm(prompt, name=name)
    return ans.strip()


def fetch_live_fact(query: str) -> str:
    try:
        q = query.lower()
        if "president" in q and ("united states" in q or "us" in q or "u.s." in q):
            r = requests.get(
                "https://en.wikipedia.org/api/rest_v1/page/summary/President_of_the_United_States",
                timeout=6,
            )
            data = r.json()
            extract = data.get("extract", "")
            sentence = extract.split(".")[0]
            return f"As of {now_str()}, {sentence}."
    except Exception as e:
        return f"Fact lookup error: {e}"
    return ""


# ==================== CODES REGISTRY =================
def _ensure_db() -> None:
    os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
    if not os.path.exists(DB_PATH):
        with open(DB_PATH, "w") as f:
            json.dump({"aliases": {}, "codes": {}}, f)


def _load_db() -> dict:
    _ensure_db()
    with open(DB_PATH, "r") as f:
        return json.load(f)


def _save_db(db: dict) -> None:
    with open(DB_PATH, "w") as f:
        json.dump(db, f, indent=2)


def _rand_code(n: int = 4) -> str:
    alphabet = string.ascii_uppercase + string.digits
    return "".join(secrets.choice(alphabet) for _ in range(n))


def save_code(alias: str, body: str) -> Tuple[str, str]:
    db = _load_db()
    code = _rand_code()
    while code in db["codes"]:
        code = _rand_code()
    db["aliases"][alias] = {"body": body, "code": code, "created": now_str()}
    db["codes"][code] = alias
    _save_db(db)
    return alias, code


def list_codes() -> List[Tuple[str, str]]:
    db = _load_db()
    return [(a, meta.get("code", "----")) for a, meta in db["aliases"].items()]


def resolve_code(key: str) -> Optional[str]:
    db = _load_db()
    k = key.strip()
    if k.startswith("#"):
        k = k[1:]
    alias = db["codes"].get(k)
    if alias:
        return db["aliases"][alias]["body"]
    if k in db["aliases"]:
        return db["aliases"][k]["body"]
    return None


def delete_code(key: str) -> bool:
    db = _load_db()
    k = key.strip()
    if k.startswith("#"):
        k = k[1:]
    alias = db["codes"].get(k)
    if alias:
        db["codes"].pop(k, None)
        db["aliases"].pop(alias, None)
        _save_db(db)
        return True
    if k in db["aliases"]:
        code = db["aliases"][k].get("code")
        if code:
            db["codes"].pop(code, None)
        db["aliases"].pop(k, None)
        _save_db(db)
        return True
    return False


# ==================== EXECUTORS =====================
def create_excel_file(
    name: str = "Quick Sheet", cols: Optional[List[str]] = None
) -> str:
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    safe = re.sub(r"[^A-Za-z0-9_\- ]", "_", name).strip().replace(" ", "_")
    path = Path(SHEETS_DIR) / f"{safe}-{ts}.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    if cols:
        ws.append(cols)
    wb.save(path)
    return str(path)


def exec_excel_command(s: str) -> str:
    m = re.search(r"name:'([^']*)'", s, flags=re.I)
    name = (m.group(1) if m else "Quick Sheet").strip()
    m = re.search(r"cols:'([^']*)'", s, flags=re.I)
    cols = [c.strip() for c in (m.group(1).split(";") if m else []) if c.strip()]
    out = create_excel_file(name=name, cols=cols)
    return f"Excel ready\nXLSX: {out}"


def make_cloud_pptx(title: str, bullets_str: str) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title or "Auto Deck"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Overview"
    tf = slide2.shapes.placeholders[1].text_frame
    tf.clear()

    bullets = [b.strip() for b in bullets_str.split(";") if b.strip()]
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            tf.add_paragraph().text = b

    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    safe = (
        re.sub(r"[^A-Za-z0-9_\- ]", "_", title or "Auto Deck").strip().replace(" ", "_")
    )
    path = CLOUD_OUT_DIR / f"{safe}-{ts}.pptx"
    prs.save(str(path))  # python-pptx expects str or file-like
    return str(path)


def exec_deck_command(s: str, to_win: bool = False) -> str:
    lower = s.strip().lower()
    if lower.startswith("ppt create "):
        title_m = re.search(r"title\s*:\s*['\"](.*?)['\"]", s, flags=re.I | re.S)
        bullets_m = re.search(r"bullets\s*:\s*['\"](.*?)['\"]", s, flags=re.I | re.S)
        title = title_m.group(1) if title_m else "Quick Deck"
        bullets_str = (
            bullets_m.group(1)
            if bullets_m
            else "Overview; Key drivers; Risks; Next steps"
        )
    else:
        text = re.sub(r"^(deck:|win:)\s*", "", s, flags=re.I).strip()
        title = "Quick Deck"
        if ";" in text:
            bullets_str = text
        else:
            ai_bullets = ask_llm("Generate 4–6 concise slide bullets about: " + text)
            cleaned = (
                ai_bullets.replace("\n", " ")
                .replace("•", " ")
                .replace(" - ", " ")
                .strip()
            )
            if ";" not in cleaned:
                parts = [p.strip(" .") for p in cleaned.split(".") if p.strip()]
                cleaned = (
                    "; ".join(parts[:6])
                    if parts
                    else "Overview; Key drivers; Risks; Next steps"
                )
            bullets_str = cleaned

    try:
        payload = {
            "command": f"ppt create title:{json.dumps(title)} bullets:{json.dumps(bullets_str)}"
        }
        headers = {"X-Signature": sign_payload(payload)}
        agent = (WIN_AGENT_URL if (to_win and WIN_AGENT_URL) else MAC_AGENT_URL).rstrip(
            "/"
        )
        url = f"{agent}/command"
        r = requests.post(url, json=payload, headers=headers, timeout=8)
        if r.ok:
            ctype = (r.headers.get("content-type") or "").lower()
            if "application/json" in ctype:
                return str(r.json().get("message", "Deck complete"))
            return (r.text or "Deck complete").strip()
    except Exception:
        path = make_cloud_pptx(title, bullets_str)
        return f"Deck complete (cloud)\nPPTX: {path}"

    path = make_cloud_pptx(title, bullets_str)
    return f"Deck complete (cloud)\nPPTX: {path}"


# ==================== ROUTING HELPERS ===============
def should_use_chat(body: str) -> bool:
    s = (body or "").strip().lower()
    if s.startswith(("deck:", "win:", "ppt create", "excel")):
        return False
    return True


# ==================== MINI TOOL PLANNER ====================
# Lightweight intent for: pull filing → make PDF → email
_TICK = re.compile(r"\b([A-Za-z]{1,5})(?:\b|:)", re.I)
_WANTS_PDF = re.compile(r"\b(pdf|send .*? pdf|email .*? pdf|as a pdf)\b", re.I)
_WANTS_EMAIL = re.compile(r"\b(send|email|mail)\b", re.I)


def _first_ticker(text: str) -> Optional[str]:
    toks = re.findall(r"\b[A-Z]{1,5}\b", text.upper())
    return toks[0] if toks else None


def try_tool_plan(body: str) -> Optional[str]:
    """
    If the user asks for latest financials/filing to be sent as a PDF,
    run edgar_pull → pdf_build → send_pdf using loaded plugins.
    Returns human message if executed, else None.
    """
    text = (body or "").strip()
    low = text.lower()
    if not any(
        k in low
        for k in (
            "10-k",
            "10q",
            "10-q",
            "financials",
            "filing",
            "latest results",
            "earnings",
        )
    ):
        return None
    if not (_WANTS_PDF.search(low) or _WANTS_EMAIL.search(low)):
        return None
    ticker = _first_ticker(text)
    if not ticker:
        return None
    needed = ["edgar_pull", "pdf_build", "send_pdf"]
    have = set(_list_plugins())
    if not all(n in have for n in needed):
        return None

    pull = _get_plugin("edgar_pull").run(ticker=ticker, prefer="10-Q")
    if not pull.get("ok"):
        return f"Could not fetch filing for {ticker}: {pull.get('error', 'unknown error')}."
    html_path = pull["html_path"]
    form = pull.get("form", "")
    fdate = pull.get("filing_date", "")

    pdf = _get_plugin("pdf_build").run(
        source_path=html_path, title=f"{ticker} {form} {fdate}"
    )
    if not pdf.get("ok"):
        return f"Fetched {ticker} {form} {fdate}, but PDF build failed: {pdf.get('error', 'unknown error')}."
    pdf_path = pdf["pdf_path"]

    to_hint = os.getenv("TO_EMAIL")
    to = to_hint or os.getenv("EMAIL_FROM") or "you@example.com"
    send = _get_plugin("send_pdf").run(
        to=to, files=[pdf_path], subject=f"{ticker} {form} {fdate}", body="Attached."
    )
    if not send.get("ok"):
        return f"PDF ready ({pdf_path}), but email failed: {send.get('error', 'unknown error')}."
    return f"Sent {ticker} {form} filed {fdate}.\nPDF: {pdf_path}"


# ==================== CORE PROCESSOR =================
def process_message(sender_id: str, body: str, channel: str = "sms") -> str:
    if os.getenv("DEBUG_LOG", "0") == "1":
        print(f"[PROC] channel={channel} sender_in={sender_id} body={body[:120]!r}")

    # Normalize SMS sender early
    if channel == "sms":
        normalized = normalize_phone(sender_id) or sender_id
        if os.getenv("DEBUG_LOG", "0") == "1":
            print(
                f"[PROC] pre-normal sms sender={sender_id} f10={last10_digits(sender_id)} → normalized={normalized}"
            )
        sender_id = normalized

    # Auth & display name
    if channel == "sms":
        f10 = last10_digits(sender_id)
        is_allowed = (
            (sender_id in ALLOWED_NUMBERS)
            or (f10 and f10 in ALLOWED_NUMBERS_LAST10)
            or _is_gv_gateway(sender_id)
        )
        if os.getenv("DEBUG_LOG", "0") == "1":
            print(f"[AUTH/SMS] sender_id={sender_id} f10={f10} allowed={is_allowed}")
        if not is_allowed:
            out = "Access denied. Unauthorized signal detected."
            return make_gv_friendly(out + jarvis_signature())
        meta = ALLOWED_NUMBERS.get(sender_id)
        if not meta and f10:
            for k, v in ALLOWED_NUMBERS.items():
                if last10_digits(k) == f10:
                    meta = v
                    break
        env_override = os.getenv("DISPLAY_NAME_CHRIS") if f10 == "5613891295" else None
        name = env_override or (meta or {}).get("name") or "Operator"
        sender_key = sender_id
    else:  # email
        addr = (parseaddr(sender_id)[1] or sender_id).strip().lower()
        if os.getenv("DEBUG_LOG", "0") == "1":
            print(f"[AUTH] email addr={addr}")
            print(f"[AUTH] EMAIL_WHITELIST={EMAIL_WHITELIST}")
            print(
                f"[AUTH] ALLOWED_EMAILS keys={[k.lower() for k in ALLOWED_EMAILS.keys()]}"
            )
        if not _is_email_allowed(addr, EMAIL_WHITELIST, ALLOWED_EMAILS, allow_gv=True):
            out = "Access denied. This email is not authorized to use Williams Cloud Control."
            return out + jarvis_signature()
        if addr in ALLOWED_EMAILS:
            name = ALLOWED_EMAILS[addr]["name"]
        else:
            name = addr.split("@")[0].replace(".", " ").title() or "Operator"
        sender_key = addr

    # Greeting (cooldown or forced)
    greeting = maybe_greeting(sender_key, name)

    def finish(msg: str) -> str:
        core = msg if msg else ""
        base = (greeting + "\n\n" + core) if greeting else core
        out = base + jarvis_signature()
        if channel == "sms":
            return make_gv_friendly(out)
        return out

    s = (body or "").strip()
    lower = s.lower()

    # If upstream sanitized the message to the fallback "ping",
    # treat it as a minimal chat request instead of a greeting.
    if lower == "ping" and os.getenv("PING_AS_FALLBACK_TO_CHAT", "1") == "1":
        s = "status report"
        lower = s.lower()

    # NEW: final SMS ingress guard (strip Subject: & admin phrasing)
    if channel == "sms":
        s = strip_subject_lines(s)
        if looks_like_admin_blob(s):
            s = "ping"
        lower = s.lower()

    # === Fast-path greetings (ALWAYS reply) ===
    if lower in {"hi", "hello", "hey", "yo"}:
        return finish("Standing by.")

    # Codes
    if lower.startswith(("code new ", "save ")):
        try:
            payload = re.sub(r"^(code new|save)\s*", "", s, flags=re.I)
            alias, raw = payload.split(":", 1)
            alias = alias.strip()
            body_for_alias = raw.strip()
            a, code = save_code(alias, body_for_alias)
            out = f"Saved code '{a}' as #{code}.\nRun with: run {a}  or  run #{code}"
            return finish(out)
        except Exception as e:
            return finish(
                f"Could not save code (format: save <alias>: <body>). Error: {e}"
            )

    if lower in ("codes", "list codes", "code list"):
        pairs = list_codes()
        if not pairs:
            return finish("No codes saved yet.")
        lines = [f"- {alias}  (#{code})" for alias, code in pairs]
        return finish("Saved codes:\n" + "\n".join(lines))

    if lower.startswith(("delete code ", "forget ")):
        key = re.sub(r"^(delete code|forget)\s*", "", s, flags=re.I).strip()
        ok = delete_code(key)
        return finish("Deleted." if ok else "Not found.")

    if lower.startswith("run "):
        key = s.split(" ", 1)[1].strip()
        stored = resolve_code(key)
        if not stored:
            return finish(f"Code '{key}' not found.")
        s = stored
        lower = s.lower()

    if lower.startswith("#"):
        stored = resolve_code(s)
        if not stored:
            return finish(f"Code '{s}' not found.")
        s = stored
        lower = s.lower()

    # Chat / search
    if should_use_chat(s):
        cleaned = re.sub(r"^(chat:|ask:)\s*", "", s, flags=re.I).strip()

        # 🔌 Try tool plan first (edgar → pdf → email)
        tool_msg = try_tool_plan(cleaned)
        if tool_msg:
            remember(cleaned, tool_msg)
            return finish(tool_msg)

        if cleaned.lower().startswith(("search ", "google ", "find ", "look up ")):
            query = re.sub(
                r"^(search|google|find|look up)\s*", "", cleaned, flags=re.I
            ).strip()
            research = research_with_perplexity(query, name=name)
            if (
                research
                and "(error" not in research.lower()
                and "missing" not in research.lower()
            ):
                brief = ask_llm(
                    "Turn the following research into a polished executive briefing. "
                    "Use concise paragraphs and elegant bullets (• or ▸), bold key headers, "
                    "end with a one-line takeaway.\n\n"
                    f"Topic: {query}\n\nResearch:\n{research}",
                    name=name,
                )
                return finish(brief)
            results = ddg_search(query, max_results=6)
            snippets = (
                format_search_results(results)
                or "No high-confidence snippets were extracted."
            )
            brief = ask_llm(
                "Using the snippets below, write a polished executive briefing on the topic. "
                "Use clean bullets (• or ▸) and a short concluding takeaway.\n\n"
                f"Topic: {query}\n\nSnippets:\n{snippets}",
                name=name,
            )
            return finish(brief)

        if cleaned.lower().startswith("scroll "):
            url = cleaned.split(" ", 1)[1].strip()
            article_summary = fetch_article(url)
            remember(f"Scroll: {url}", article_summary)
            return finish(article_summary)

        live = fetch_live_fact(cleaned)
        if live:
            remember(cleaned, live)
            return finish(live)

        if needs_freshness(cleaned):
            # Hint the router + set the time context
            cleaned = (
                "use: perplexity\n"
                f"{cleaned}\n\n"
                f"Context: Answer as of {now_str()} (current year: {datetime.now(ZoneInfo(TZ)).year}). "
                "Cite sources inline."
            )

        ai_reply = ask_llm(cleaned or "Status report.", name=name)
        remember(cleaned, ai_reply)
        return finish(ai_reply)

    # Excel
    if lower.startswith("excel"):
        try:
            msg = exec_excel_command(s)
            remember(s, msg)
            return finish(msg)
        except Exception as e:
            err = f"{e}\n{traceback.format_exc(limit=1)}"
            return finish(f"Excel exception: {err}")

    # Deck
    to_win = lower.startswith("win:")
    if lower.startswith(("deck:", "win:", "ppt create ")):
        try:
            msg = exec_deck_command(s, to_win=to_win)
            remember(s, msg)
            return finish(msg)
        except Exception as e:
            err = f"{e}\n{traceback.format_exc(limit=1)}"
            return finish(f"Agent exception: {err}")

    # ===== SEC / EDGAR COMMANDS =====
    # Examples:
    #   "sec aapl"         -> latest 10-K, extract Income Statement to XLSX
    #   "edgar tsla 10-k"  -> same as above
    #   "filing msft 10k"  -> same as above
    if any(x in lower for x in ["sec", "edgar", "10-k", "10k", "filing"]):
        # parse ticker: first ALLCAPS/alnum token after the keyword, fallback to last word
        m = re.search(
            r"(?:sec|edgar|filing|10-?k)\s+([A-Za-z0-9\.\-]{1,10})", lower, flags=re.I
        )
        ticker = None
        if m:
            ticker = m.group(1).upper()
        else:
            toks = [t for t in re.split(r"\s+", s.strip()) if t]
            if toks:
                cand = toks[-1]
                if re.fullmatch(r"[A-Za-z0-9\.\-]{1,10}", cand):
                    ticker = cand.upper()

        if not ticker:
            return finish("Usage: sec <TICKER>  (e.g., sec AAPL)")

        ok, msg = sec_pull_10k_income_to_xlsx(ticker)
        if not ok:
            return finish(f"SEC error ({ticker}): {msg}")

        # Return a line that your GV path turns into a link automatically
        return finish(f"Income statement extracted.\nXLSX: {msg}")

    # Fallback
    return finish("Command not recognized.")


# ==================== HEALTH ========================
@app.get("/health")
def health() -> Dict[str, str]:
    return {"ok": True, "time": now_str()}


# ==================== TOOLS: list & invoke ==========
@app.get("/tools")
def list_tools() -> Dict[str, Any]:
    return {"ok": True, "plugins": _list_plugins()}


@app.post("/tool")
async def run_tool(req: Request):
    data = await req.json()
    tool = data.get("tool")
    args = data.get("args", {}) or {}
    try:
        mod = _get_plugin(tool)
    except Exception:
        return JSONResponse(
            {
                "ok": False,
                "error": f"unknown tool '{tool}'",
                "available": _list_plugins(),
            },
            status_code=400,
        )
    try:
        out = mod.run(**args)
        return JSONResponse(out)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)


# ==================== TWILIO SMS (optional path) ====
# (Not required for GV/Gmail flow; kept for local sanity checks)
@app.post("/twilio/sms")
async def sms_in(req: Request):
    try:
        form = dict(await req.form())
        sender = (form.get("From") or "").strip()
        body = (form.get("Body") or "").strip()
    except Exception:
        return twiml("Malformed request.")
    sender = normalize_phone(sender) or sender
    body_clean = extract_gv_sms_payload(body)
    # NEW: kill Subject: and admin blobs
    body_clean = strip_subject_lines(body_clean)
    if looks_like_admin_blob(body_clean):
        body_clean = "ping"
    if not body_clean.strip():
        if os.getenv("DEBUG_LOG", "0") == "1":
            print("[TWILIO] Cleaned payload empty; using 'ping'")
        body_clean = "ping"
    text_reply = process_message(sender, body_clean, channel="sms")
    text_reply = make_gv_friendly(text_reply)
    return twiml(text_reply)


# ==================== EMAIL INBOUND (GV Gateway) ====
@app.post("/email/inbound")
async def email_inbound(req: Request):
    ctype = (req.headers.get("content-type") or "").lower()
    sender = subject = text = ""
    if "application/json" in ctype:
        data = await req.json()
        sender = (data.get("from") or data.get("sender") or "").strip()
        subject = (data.get("subject") or "").strip()
        text = (
            data.get("text") or data.get("body") or data.get("stripped-text") or ""
        ).strip()
    else:
        form = dict(await req.form())
        sender = (form.get("from") or form.get("sender") or "").strip()
        subject = (form.get("subject") or "").strip()
        text = (
            form.get("text")
            or form.get("stripped-text")
            or form.get("body-plain")
            or ""
        ).strip()

    if os.getenv("DEBUG_LOG", "0") == "1":
        print(
            f"[INBOUND] ctype={ctype} sender={sender} subject={subject!r} len(text)={len(text or '')}"
        )

    if not sender:
        raise HTTPException(status_code=400, detail="missing sender")

    command_text = text if text else subject
    sender_addr = (parseaddr(sender)[1] or sender).lower()

    # Treat GV gateway as SMS so we bypass the email allowlist
    if _is_gv_gateway(sender_addr):
        blob = f"{subject or ''} {text or ''}".lower()
        if (
            os.getenv("GV_HARD_IGNORE_SETTINGS", "1") == "1"
            and "voice.google.com/settings#messaging" in blob
        ):
            if os.getenv("DEBUG_LOG", "0") == "1":
                print("[GV] Hard-ignore GV settings notice in /email/inbound")
            return {"ok": True, "gv": True, "ignored": "admin_notice"}

        # heuristic ignore (optional)
        if is_gv_admin_notice(subject, text, headers={"from": sender_addr}):
            if os.getenv("DEBUG_LOG", "0") == "1":
                print("[GV] Admin notice detected — ignored.")
            return {"ok": True, "gv": True, "ignored": "admin_notice"}

        local = sender_addr.split("@", 1)[0]
        phone = (
            normalize_phone(local) or normalize_phone(subject) or normalize_phone(text)
        )
        sms_sender = normalize_phone(phone or sender_addr) or (phone or sender_addr)

        clean_text = extract_gv_sms_payload(command_text)
        # NEW: kill Subject: and admin blobs
        clean_text = strip_subject_lines(clean_text)
        if looks_like_admin_blob(clean_text):
            clean_text = "ping"

        if not clean_text.strip():
            if os.getenv("DEBUG_LOG", "0") == "1":
                print("[GV] Cleaned payload was empty; falling back to 'ping'")
            clean_text = "ping"

        if os.getenv("DEBUG_LOG", "0") == "1":
            print(f"[GV] Processing SMS from {sms_sender} … {clean_text[:80]!r}")

        reply = process_message(sms_sender, clean_text, channel="sms")
        if os.getenv("DEBUG_LOG", "0") == "1":
            print(f"[GV] sms_sender={sms_sender} reply_len={len(reply)}")

        attachments = _extract_paths(reply)
        reply = _replace_paths_with_links(reply, attachments)
        reply = make_gv_friendly(reply)

        # Deliver: Gmail API → SMTP → fallback file
        delivered = False
        try:
            svc = _gmail_service()
            _gmail_send_plain(svc, sender_addr, f"Re: {subject or 'Command'}", reply)
            delivered = True
            if os.getenv("DEBUG_LOG", "0") == "1":
                print("[GV] Sent via Gmail API")
        except Exception as e:
            if os.getenv("DEBUG_LOG", "0") == "1":
                print("[GV] Gmail API send failed:", e)

        if not delivered:
            res = send_email(
                sender_addr, f"Re: {subject or 'Command'}", reply, attachments=None
            )
            if os.getenv("DEBUG_LOG", "0") == "1":
                print("[GV] SMTP result:", res)
            delivered = not res.startswith("(Email disabled")

        if not delivered:
            log_path = Path(CLOUD_OUT_DIR) / f"gv_out_{int(time.time())}.txt"
            log_path.write_text(
                f"TO: {sender_addr}\nSUBJECT: Re: {subject or 'Command'}\n\n{reply}"
            )
            if os.getenv("DEBUG_LOG", "0") == "1":
                print(f"[GV] Wrote fallback file: {log_path}")

        return {"ok": True, "gv": True}

    # Normal email path
    reply2 = process_message(sender_addr, command_text, channel="email")
    attachments2 = _extract_paths(reply2)
    sent = send_email(
        sender_addr,
        f"[Williams Cloud] Re: {subject or 'Command'}",
        reply2,
        attachments=attachments2,
    )
    return {"ok": True, "sent": sent}


# ==================== Gmail API helpers (GV bridge) ==
def _gmail_service():
    if not _GOOGLE_OK:
        raise RuntimeError(
            "Google libs not installed. pip install google-api-python-client google-auth-httplib2 google-auth-oauthlib"
        )
    if not os.path.exists(GV_TOKEN_JSON) or not os.path.exists(GV_CLIENT_JSON):
        raise RuntimeError("Missing token.json or client_secret.json for Gmail.")
    creds = Credentials.from_authorized_user_file(GV_TOKEN_JSON, GMAIL_SCOPES)
    if not creds.valid:
        if creds.expired and creds.refresh_token:
            creds.refresh(GoogleRequest())
            with open(GV_TOKEN_JSON, "w") as f:
                f.write(creds.to_json())
        else:
            raise RuntimeError("Gmail token invalid; re-authorize.")
    return build("gmail", "v1", credentials=creds)


def _gmail_list_unread_gv(svc) -> List[Dict[str, Any]]:
    res = svc.users().messages().list(userId="me", q=GMAIL_QUERY).execute()
    return res.get("messages", []) or []


def _gmail_get_plaintext(
    svc, msg_id: str
) -> Tuple[dict, Dict[str, str], str, str, str]:
    m = svc.users().messages().get(userId="me", id=msg_id, format="full").execute()
    headers = {
        h["name"].lower(): h["value"] for h in m.get("payload", {}).get("headers", [])
    }
    frm = headers.get("from", "")
    subj = headers.get("subject", "")

    def _dec(b64: str) -> str:
        return base64.urlsafe_b64decode(b64.encode()).decode("utf-8", "ignore")

    body = ""
    payload = m.get("payload", {})
    if "data" in payload.get("body", {}):
        body = _dec(payload["body"]["data"])
    else:
        for part in payload.get("parts", []):
            if part.get("mimeType") == "text/plain" and "data" in part.get("body", {}):
                body = _dec(part["body"]["data"])
                break
    return m, headers, frm, subj, (body or "").strip()


def _gmail_mark_read(svc, msg_id: str) -> None:
    svc.users().messages().modify(
        userId="me", id=msg_id, body={"removeLabelIds": ["UNREAD"]}
    ).execute()


def _gmail_send_plain(svc, to_addr: str, subject: str, body: str) -> None:
    msg = MIMEText(body)
    msg["To"] = to_addr
    msg["Subject"] = subject or "Re:"
    raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
    svc.users().messages().send(userId="me", body={"raw": raw}).execute()


def _gmail_send_with_attachments(
    svc, to_addr: str, subject: str, body: str, attachments: Optional[List[str]] = None
) -> None:
    msg = MIMEMultipart()
    msg["To"] = to_addr
    msg["Subject"] = subject or "Re:"
    msg.attach(MIMEText(body, "plain"))
    for path in attachments or []:
        if not os.path.isfile(path):
            continue
        ctype, _ = mimetypes.guess_type(path)
        maintype, subtype = (ctype or "application/octet-stream").split("/", 1)
        with open(path, "rb") as f:
            part = MIMEBase(maintype, subtype)
            part.set_payload(f.read())
        encoders.encode_base64(part)
        part.add_header(
            "Content-Disposition", "attachment", filename=os.path.basename(path)
        )
        msg.attach(part)
    raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
    svc.users().messages().send(userId="me", body={"raw": raw}).execute()


# ==================== GV POLLING LOOP =================
def run_gmail_polling_once() -> None:
    """
    Reads unread GV-SMS emails, normalizes sender phone, routes through process_message(),
    emails reply back (GV converts to SMS), marks message read.
    """
    svc = _gmail_service()
    msgs = _gmail_list_unread_gv(svc)
    if os.getenv("DEBUG_LOG", "0") == "1":
        print(f"[GV] Unread GV messages: {len(msgs)}")
    for m in msgs:
        full, headers, frm, subj, text = _gmail_get_plaintext(svc, m["id"])

        if os.getenv("DEBUG_LOG", "0") == "1":
            print(f"[GV] Seen: subj={subj!r} from={frm!r} len(text)={len(text or '')}")

        # HARD-IGNORE exact GV settings notification (env-gated)
        if (
            os.getenv("GV_HARD_IGNORE_SETTINGS", "1") == "1"
            and "voice.google.com/settings#messaging"
            in f"{subj or ''} {text or ''}".lower()
        ):
            if os.getenv("DEBUG_LOG", "0") == "1":
                print(
                    "[GV] Hard-ignore GV settings notice in poller — marked read, no reply."
                )
            _gmail_mark_read(svc, m["id"])
            continue

        # Heuristic ignore
        if is_gv_admin_notice(subj, text, headers):
            if os.getenv("DEBUG_LOG", "0") == "1":
                print("[GV] Admin notice detected in poller — marked read, no reply.")
            _gmail_mark_read(svc, m["id"])
            continue

        phone_e164 = extract_phone_from_gv_headers(headers, subj, text)
        reply_to_hdr = headers.get("reply-to") or headers.get("from") or ""
        to_addr = reply_to_hdr

        if _is_gv_gateway(to_addr):
            sms_sender = normalize_phone(phone_e164 or to_addr) or (
                phone_e164 or to_addr
            )
            incoming = text or subj or ""
            clean_incoming = extract_gv_sms_payload(incoming)
            # NEW: kill Subject: and admin blobs
            clean_incoming = strip_subject_lines(clean_incoming)
            if looks_like_admin_blob(clean_incoming):
                clean_incoming = "ping"
            if not clean_incoming.strip():
                # keep original text instead of forcing 'ping'
                fallback_blob = (text or subj or "").strip()
                clean_incoming = fallback_blob if fallback_blob else "[empty]"

            if os.getenv("DEBUG_LOG", "0") == "1":
                print(
                    f"[GV] Processing SMS from {sms_sender} … {clean_incoming[:80]!r}"
                )

            reply_text = process_message(sms_sender, clean_incoming, channel="sms")
        else:
            reply_text = process_message(to_addr, text or subj or "", channel="email")

        attachments = _extract_paths(reply_text)
        if _is_gv_gateway(to_addr):
            reply_text = _replace_paths_with_links(reply_text, attachments)
            reply_text = make_gv_friendly(reply_text)

            delivered = False
            try:
                _gmail_send_plain(svc, to_addr, subj, reply_text)
                delivered = True
                if os.getenv("DEBUG_LOG", "0") == "1":
                    print("[GV/Poller] Sent via Gmail API")
            except Exception as e:
                if os.getenv("DEBUG_LOG", "0") == "1":
                    print("[GV/Poller] Gmail API send failed:", e)

            if not delivered:
                res = send_email(to_addr, subj or "Re:", reply_text, attachments=None)
                if os.getenv("DEBUG_LOG", "0") == "1":
                    print("[GV/Poller] SMTP result:", res)
                delivered = not res.startswith("(Email disabled")

            if not delivered:
                log_path = Path(CLOUD_OUT_DIR) / f"gv_out_{int(time.time())}.txt"
                log_path.write_text(
                    f"TO: {to_addr}\nSUBJECT: {subj or 'Re:'}\n\n{reply_text}"
                )
                if os.getenv("DEBUG_LOG", "0") == "1":
                    print(f"[GV/Poller] Wrote fallback file: {log_path}")
        else:
            _gmail_send_with_attachments(svc, to_addr, subj, reply_text, attachments)

        _gmail_mark_read(svc, m["id"])


def run_gmail_polling_loop() -> None:
    print(f"[GV] Polling Gmail every {POLL_SECONDS}s (query: {GMAIL_QUERY})")
    while True:
        try:
            run_gmail_polling_once()
        except Exception as e:
            print("Gmail polling error:", e)
        time.sleep(POLL_SECONDS)


# ==================== DEBUG: Gmail account whoami ===
@app.get("/debug/gmail/whoami")
def gmail_whoami() -> Dict[str, Any]:
    try:
        svc = _gmail_service()
        prof = svc.users().getProfile(userId="me").execute()
        return {"ok": True, "emailAddress": prof.get("emailAddress")}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ==================== DEBUG: trigger one poll now ===
@app.post("/debug/gmail_once")
def gmail_once() -> Dict[str, Any]:
    try:
        run_gmail_polling_once()
        return {"ok": True}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ==================== DEBUG: GV simulator =============
@app.post("/debug/simulate_gv")
async def debug_simulate_gv(
    sender_local: Optional[str] = None,
    subject: Optional[str] = None,
    text: Optional[str] = None,
    # Form fallbacks so `curl -d` works
    sender_local_form: Optional[str] = Form(default=None),
    subject_form: Optional[str] = Form(default=None),
    text_form: Optional[str] = Form(default=None),
    payload: Optional[dict] = Body(default=None),
):
    # unify inputs from query / form / JSON
    if payload:
        sender_local = sender_local or payload.get("sender_local")
        subject = subject or payload.get("subject")
        text = text or payload.get("text")
    sender_local = sender_local or sender_local_form
    subject = (subject or subject_form or "") or ""
    text = (text or text_form or "") or ""

    if not sender_local:
        raise HTTPException(status_code=400, detail="sender_local is required")

    # Ensure GV domain suffix
    sender_addr = sender_local.lower()
    if not sender_addr.endswith("@txt.voice.google.com"):
        sender_addr += "@txt.voice.google.com"

    # Match poller extraction priority
    headers = {
        "from": sender_addr,
        "reply-to": sender_addr,
        "to": sender_addr,
        "delivered-to": sender_addr,
    }
    phone = extract_phone_from_gv_headers(headers, subject, text)
    if not phone:
        phone = (
            normalize_phone(subject)
            or normalize_phone(text)
            or normalize_phone(sender_addr.split("@", 1)[0])
            or sender_addr
        )

    incoming = text or subject or "ping"
    clean_incoming = extract_gv_sms_payload(incoming) or "ping"
    # NEW: kill Subject: and admin blobs in simulator too
    clean_incoming = strip_subject_lines(clean_incoming)
    if looks_like_admin_blob(clean_incoming):
        clean_incoming = "ping"

    reply = process_message(str(phone), clean_incoming, channel="sms")
    reply = make_gv_friendly(_replace_paths_with_links(reply, _extract_paths(reply)))

    return {
        "ok": True,
        "sms_sender": phone,
        "reply_preview": reply[:900],
        "note": "Simulator uses poller-matching phone extraction + allowlist.",
    }


# ==================== DEBUG: trigger self-improver ===
@app.post("/debug/self_improve")
def debug_self_improve() -> Dict[str, Any]:
    """
    Runs the improver once (dry-run first); requires tools/auto_improver.py present.
    """
    try:
        logs: List[str] = []
        r1 = subprocess.run(
            [sys.executable, "tools/auto_improver.py", "--dry-run", "--show"],
            capture_output=True,
            text=True,
            timeout=600,
        )
        logs.append(r1.stdout[-4000:])
        if r1.returncode == 0:
            r2 = subprocess.run(
                [sys.executable, "tools/auto_improver.py", "--apply"],
                capture_output=True,
                text=True,
                timeout=1200,
            )
            logs.append(r2.stdout[-4000:])
            ok = r2.returncode == 0
        else:
            ok = False
        return {"ok": ok, "log_tail": "\n\n---\n\n".join(logs)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ==================== MAIN (local dev helper) ========
if __name__ == "__main__":
    # Run FastAPI + GV poller together for local dev
    import uvicorn

    t = threading.Thread(target=run_gmail_polling_loop, daemon=True)
    t.start()
    uvicorn.run(app, host="0.0.0.0", port=8000)

# ==================== SEC / EDGAR HELPERS ====================
SEC_BASE = "https://www.sec.gov/"
SEC_DATA = "https://data.sec.gov/"
COMPANY_TICKERS_URL = "https://www.sec.gov/files/company_tickers.json"


def _sec_headers() -> Dict[str, str]:
    appname = os.getenv("SEC_APP_NAME", "Williams Cloud Control")
    email = os.getenv("SEC_CONTACT_EMAIL", "example@example.com")
    return {
        "User-Agent": f"{appname} ({email})",
        "Accept": "text/html,application/json;q=0.9,*/*;q=0.8",
    }


def _sec_get_json(url: str, timeout: int = 12) -> dict:
    r = requests.get(url, headers=_sec_headers(), timeout=timeout)
    r.raise_for_status()
    return r.json()


def _sec_get_text(url: str, timeout: int = 15) -> str:
    r = requests.get(url, headers=_sec_headers(), timeout=timeout)
    r.raise_for_status()
    return r.text


# Cache ticker map in-process to avoid repeated fetches
_TICKER_MAP: Optional[Dict[str, str]] = None


def _get_ticker_map() -> Dict[str, str]:
    global _TICKER_MAP
    if _TICKER_MAP is not None:
        return _TICKER_MAP
    try:
        data = _sec_get_json(COMPANY_TICKERS_URL)
        # File is { "0": {"cik_str": 320193, "ticker":"AAPL","title":"Apple Inc."}, ...}
        mp: Dict[str, str] = {}
        for _, row in data.items():
            t = (row.get("ticker") or "").upper().strip()
            if not t:
                continue
            cik = str(row.get("cik_str") or "").strip()
            if cik:
                mp[t] = cik.zfill(10)
        _TICKER_MAP = mp
        return mp
    except Exception:
        _TICKER_MAP = {}
        return _TICKER_MAP


def _cik_from_ticker(ticker: str) -> Optional[str]:
    t = (ticker or "").upper().strip()
    if not t:
        return None
    mp = _get_ticker_map()
    cik = mp.get(t)
    return cik


def _latest_10k_meta(cik: str) -> Optional[dict]:
    """
    Use submissions API to find latest non-amended 10-K.
    """
    url = urljoin(SEC_DATA, f"submissions/CIK{cik}.json")
    try:
        js = _sec_get_json(url)
        recent = js.get("filings", {}).get("recent", {})
        forms = recent.get("form", []) or []
        accs = recent.get("accessionNumber", []) or []
        prims = recent.get("primaryDocument", []) or []
        dates = recent.get("filingDate", []) or []
        for i, form in enumerate(forms):
            if (form or "").upper() == "10-K":  # skip 10-K/A
                acc = (accs[i] or "").replace("-", "")
                prim = prims[i]
                fd = dates[i]
                if acc and prim:
                    return {"accession_nodash": acc, "primary": prim, "filingDate": fd}
    except Exception:
        pass
    return None


def _build_primary_doc_url(cik: str, accession_nodash: str, primary: str) -> str:
    # https://www.sec.gov/Archives/edgar/data/{cik}/{accession_nodash}/{primary}
    return urljoin(
        SEC_BASE, f"Archives/edgar/data/{int(cik)}/{accession_nodash}/{primary}"
    )


def _extract_best_income_statement(html_text: str) -> Optional[pd.DataFrame]:
    """
    Parse all tables, score by income-statement keywords, return the best match.
    """
    try:
        # Fast path: pandas read_html
        tables = pd.read_html(html_text, flavor="lxml")
    except Exception:
        # Fallback via BeautifulSoup -> then pandas on each <table>
        soup = BeautifulSoup(html_text, "lxml")
        tables = []
        for tbl in soup.find_all("table"):
            try:
                df = pd.read_html(str(tbl))[0]
                tables.append(df)
            except Exception:
                continue
    if not tables:
        return None

    KEYWORDS = [
        "net income",
        "net loss",
        "total revenue",
        "net sales",
        "operating income",
        "gross profit",
        "cost of sales",
        "earnings per share",
    ]

    def score_df(df: pd.DataFrame) -> int:
        try:
            txt = " ".join(
                [str(x) for x in df.columns.tolist()]
                + [str(x) for x in df.head(10).to_numpy().ravel().tolist()]
            ).lower()
        except Exception:
            txt = str(df.head(10)).lower()
        score = sum(int(k in txt) for k in KEYWORDS)
        # prefer wider & taller tables modestly
        rows, cols = df.shape
        score += min(rows, 50) // 10 + min(cols, 12) // 4
        return score

    ranked = sorted(tables, key=score_df, reverse=True)
    best = ranked[0]
    # Clean up: drop all-empty columns/rows
    best = best.dropna(how="all", axis=0).dropna(how="all", axis=1)
    return best


def _save_income_statement_to_xlsx(
    ticker: str, filing_date: str, df: pd.DataFrame
) -> str:
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    safe_t = re.sub(r"[^A-Za-z0-9_\-]", "_", ticker.upper())
    filename = f"{safe_t}_10K_Income_{filing_date}_{ts}.xlsx"
    path = CLOUD_OUT_DIR / filename
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Income Statement")
    return str(path)


def sec_pull_10k_income_to_xlsx(ticker: str) -> Tuple[bool, str]:
    """
    Returns (ok, message_or_path). On success, message is XLSX path.
    """
    t = (ticker or "").upper().strip()
    if not t:
        return False, "Ticker missing."
    cik = _cik_from_ticker(t)
    if not cik:
        return False, f"Could not resolve CIK for {t}."
    meta = _latest_10k_meta(cik)
    if not meta:
        return False, f"No recent 10-K found for CIK {cik}."
    url = _build_primary_doc_url(cik, meta["accession_nodash"], meta["primary"])
    try:
        html = _sec_get_text(url)
    except Exception as e:
        return False, f"Fetch error: {e}"
    df = _extract_best_income_statement(html)
    if df is None or df.empty:
        return (
            False,
            "Could not locate an income statement table in the primary document.",
        )
    xlsx_path = _save_income_statement_to_xlsx(t, meta.get("filingDate", ""), df)
    return True, xlsx_path
