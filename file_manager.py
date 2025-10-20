#!/usr/bin/env python3
"""
Enhanced File Management System
Handles PPT creation, editing, PDF generation, and file operations.
"""

import os
import re
import logging
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
import asyncio

# Add project root to path
import sys
sys.path.insert(0, str(Path(__file__).parent))

from config.settings_simple import get_settings
from agent_v2 import superchat
from cloud_storage import storage_adapter

logger = logging.getLogger(__name__)

class FileManager:
    """Manages file operations for the AI system."""
    
    def __init__(self):
        self.settings = get_settings()
        self.cloud_out_dir = self.settings.cloud_out_dir
        self.sheets_dir = self.settings.sheets_dir
        self.public_base_url = self.settings.public_base_url
        
        # Ensure directories exist
        self.cloud_out_dir.mkdir(parents=True, exist_ok=True)
        self.sheets_dir.mkdir(parents=True, exist_ok=True)
        
        # File type handlers (placeholder methods)
        self.handlers = {
            'pptx': self._handle_pptx,
            'pdf': self._handle_pdf,
            'xlsx': self._handle_xlsx,
            'docx': self._handle_docx,
            'txt': self._handle_txt
        }
    
    def _handle_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Handle PPTX files."""
        return {"type": "pptx", "path": str(file_path)}
    
    def _handle_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Handle PDF files."""
        return {"type": "pdf", "path": str(file_path)}
    
    def _handle_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """Handle XLSX files."""
        return {"type": "xlsx", "path": str(file_path)}
    
    def _handle_docx(self, file_path: Path) -> Dict[str, Any]:
        """Handle DOCX files."""
        return {"type": "docx", "path": str(file_path)}
    
    def _handle_txt(self, file_path: Path) -> Dict[str, Any]:
        """Handle TXT files."""
        return {"type": "txt", "path": str(file_path)}
    
    async def process_file_request(self, request: str, user_id: str = None) -> Dict[str, Any]:
        """Process a file-related request from the user."""
        request_lower = request.lower()
        
        # Detect file operation type
        if any(keyword in request_lower for keyword in ['create ppt', 'make presentation', 'powerpoint']):
            return await self._create_presentation(request, user_id)
        
        elif any(keyword in request_lower for keyword in ['edit ppt', 'modify presentation', 'update powerpoint']):
            return await self._edit_presentation(request, user_id)
        
        elif any(keyword in request_lower for keyword in ['10-k', '10k', 'sec filing', 'edgar']):
            return await self._handle_sec_filing(request, user_id)
        
        elif any(keyword in request_lower for keyword in ['send pdf', 'email pdf', 'mail pdf']):
            return await self._send_pdf(request, user_id)
        
        elif any(keyword in request_lower for keyword in ['create excel', 'make spreadsheet']):
            return await self._create_excel(request, user_id)
        
        else:
            return {
                "success": False,
                "message": "I didn't understand the file operation. Try: 'create ppt', 'edit presentation', 'pull Apple 10-K', or 'send pdf'"
            }
    
    async def _create_presentation(self, request: str, user_id: str = None) -> Dict[str, Any]:
        """Create a PowerPoint presentation."""
        try:
            # Use AI to generate presentation content
            ai_prompt = f"""
Create a PowerPoint presentation based on this request: {request}

Please provide:
1. Title slide content
2. 3-5 main bullet points for the presentation
3. Any specific data or information to include

Format your response as:
TITLE: [presentation title]
BULLETS:
• [bullet point 1]
• [bullet point 2]
• [bullet point 3]
• [bullet point 4]
• [bullet point 5]
"""
            
            ai_response = await superchat(ai_prompt)
            
            # Parse AI response
            title_match = re.search(r'TITLE:\s*(.+)', ai_response)
            bullets_match = re.search(r'BULLETS:(.*?)(?=\n[A-Z]+:|$)', ai_response, re.DOTALL)
            
            title = title_match.group(1).strip() if title_match else "AI Generated Presentation"
            bullets_text = bullets_match.group(1).strip() if bullets_match else "• AI generated content"
            
            # Create the presentation
            from pptx import Presentation
            prs = Presentation()
            
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Content slide
            content_slide = prs.slides.add_slide(prs.slide_layouts[1])
            content_title = content_slide.shapes.title
            content_title.text = "Key Points"
            
            content_body = content_slide.placeholders[1]
            content_body.text = bullets_text
            
            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
            safe_title = re.sub(r"[^A-Za-z0-9_\- ]", "_", title).strip().replace(" ", "_")
            # Limit filename length to avoid filesystem issues
            safe_title = safe_title[:50]  # Limit to 50 characters
            filename = f"{safe_title}-{timestamp}.pptx"
            
            # Save to temporary file first
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                prs.save(tmp_file.name)
                with open(tmp_file.name, 'rb') as f:
                    file_content = f.read()
                os.unlink(tmp_file.name)
            
            # Upload to cloud storage
            storage_result = await storage_adapter.save_file(filename, file_content, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            
            if not storage_result.get("success"):
                return {
                    "success": False,
                    "message": f"❌ Error saving presentation: {storage_result.get('error', 'Unknown error')}"
                }
            
            public_url = storage_result["public_url"]
            
            return {
                "success": True,
                "message": f"✅ Presentation created: {title}",
                "file_path": storage_result.get("storage_path", filename),
                "public_url": public_url,
                "filename": filename,
                "title": title,
                "bullets": bullets_text
            }
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            return {
                "success": False,
                "message": f"❌ Error creating presentation: {str(e)}"
            }
    
    async def _edit_presentation(self, request: str, user_id: str = None) -> Dict[str, Any]:
        """Edit an existing PowerPoint presentation."""
        try:
            # Extract file path or name from request
            file_match = re.search(r'(?:file|presentation|ppt)[:\s]+([^\s]+)', request, re.I)
            if not file_match:
                return {
                    "success": False,
                    "message": "Please specify which presentation to edit. Example: 'edit presentation myfile.pptx'"
                }
            
            filename = file_match.group(1)
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            file_path = self.cloud_out_dir / filename
            if not file_path.exists():
                return {
                    "success": False,
                    "message": f"❌ File not found: {filename}"
                }
            
            # Load existing presentation
            from pptx import Presentation
            prs = Presentation(str(file_path))
            
            # Use AI to determine edits
            ai_prompt = f"""
I have a PowerPoint presentation file: {filename}
User wants to edit it with this request: {request}

Please provide specific editing instructions:
1. What slides to modify
2. What content to add/remove/change
3. Any new slides to add

Format as:
EDIT_SLIDE_1: [instructions for slide 1]
EDIT_SLIDE_2: [instructions for slide 2]
ADD_SLIDE: [content for new slide]
"""
            
            ai_response = await superchat(ai_prompt)
            
            # Apply edits (simplified - in real implementation, you'd parse and apply each edit)
            # For now, add a new slide with the edit request
            new_slide = prs.slides.add_slide(prs.slide_layouts[1])
            new_slide.shapes.title.text = "AI Edits"
            new_slide.placeholders[1].text = f"Edit request: {request}\n\nAI suggestions: {ai_response[:200]}..."
            
            # Save edited presentation
            timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
            edited_filename = f"edited_{filename.replace('.pptx', '')}-{timestamp}.pptx"
            edited_path = self.cloud_out_dir / edited_filename
            prs.save(str(edited_path))
            
            public_url = f"{self.public_base_url}/files/{edited_filename}"
            
            return {
                "success": True,
                "message": f"✅ Presentation edited: {edited_filename}",
                "file_path": str(edited_path),
                "public_url": public_url,
                "original_file": str(file_path),
                "edits_applied": ai_response[:200] + "..."
            }
            
        except Exception as e:
            logger.error(f"Error editing presentation: {e}")
            return {
                "success": False,
                "message": f"❌ Error editing presentation: {str(e)}"
            }
    
    async def _handle_sec_filing(self, request: str, user_id: str = None) -> Dict[str, Any]:
        """Handle SEC filing requests (like Apple 10-K)."""
        try:
            # Extract ticker symbol
            ticker_match = re.search(r'\b([A-Z]{1,5})\b', request.upper())
            if not ticker_match:
                return {
                    "success": False,
                    "message": "Please specify a ticker symbol. Example: 'pull Apple 10-K' or 'get AAPL 10-K'"
                }
            
            ticker = ticker_match.group(1)
            
            # Use existing tool plan functionality
            from cloud import try_tool_plan
            result = try_tool_plan(f"pull {ticker} 10-K")
            
            if result:
                return {
                    "success": True,
                    "message": f"✅ {ticker} 10-K processed and sent via email",
                    "details": result
                }
            else:
                return {
                    "success": False,
                    "message": f"❌ Could not process {ticker} 10-K. Make sure the edgar_pull, pdf_build, and send_pdf plugins are available."
                }
                
        except Exception as e:
            logger.error(f"Error handling SEC filing: {e}")
            return {
                "success": False,
                "message": f"❌ Error processing SEC filing: {str(e)}"
            }
    
    async def _send_pdf(self, request: str, user_id: str = None) -> Dict[str, Any]:
        """Send PDF via email."""
        try:
            # Extract file path from request
            file_match = re.search(r'(?:file|pdf)[:\s]+([^\s]+)', request, re.I)
            if not file_match:
                return {
                    "success": False,
                    "message": "Please specify which PDF to send. Example: 'send pdf myfile.pdf'"
                }
            
            filename = file_match.group(1)
            if not filename.endswith('.pdf'):
                filename += '.pdf'
            
            file_path = self.cloud_out_dir / filename
            if not file_path.exists():
                return {
                    "success": False,
                    "message": f"❌ PDF file not found: {filename}"
                }
            
            # Use existing email functionality
            from cloud import send_email
            to_email = self.settings.email_from or "you@example.com"
            
            result = send_email(
                to_email,
                f"PDF Document: {filename}",
                f"Please find attached: {filename}",
                [str(file_path)]
            )
            
            if result.get("ok"):
                return {
                    "success": True,
                    "message": f"✅ PDF sent to {to_email}",
                    "file_path": str(file_path),
                    "recipient": to_email
                }
            else:
                return {
                    "success": False,
                    "message": f"❌ Error sending PDF: {result.get('error', 'Unknown error')}"
                }
                
        except Exception as e:
            logger.error(f"Error sending PDF: {e}")
            return {
                "success": False,
                "message": f"❌ Error sending PDF: {str(e)}"
            }
    
    async def _create_excel(self, request: str, user_id: str = None) -> Dict[str, Any]:
        """Create Excel spreadsheet."""
        try:
            # Use AI to generate spreadsheet content
            ai_prompt = f"""
Create an Excel spreadsheet based on this request: {request}

Please provide:
1. Sheet name
2. Column headers
3. Sample data (3-5 rows)

Format as:
SHEET_NAME: [name]
COLUMNS: [col1, col2, col3]
DATA:
[data1, data2, data3]
[data4, data5, data6]
[data7, data8, data9]
"""
            
            ai_response = await superchat(ai_prompt)
            
            # Parse AI response
            sheet_match = re.search(r'SHEET_NAME:\s*(.+)', ai_response)
            cols_match = re.search(r'COLUMNS:\s*(.+)', ai_response)
            data_match = re.search(r'DATA:(.*?)(?=\n[A-Z]+:|$)', ai_response, re.DOTALL)
            
            sheet_name = sheet_match.group(1).strip() if sheet_match else "AI Generated Sheet"
            # Clean sheet name for Excel compatibility
            sheet_name = re.sub(r'[^\w\s-]', '', sheet_name)[:31]  # Excel sheet name limit
            columns = [col.strip() for col in cols_match.group(1).split(',')] if cols_match else ["Column 1", "Column 2", "Column 3"]
            
            # Create Excel file
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            # Add headers
            for i, col in enumerate(columns, 1):
                ws.cell(row=1, column=i, value=col)
            
            # Add data if provided
            if data_match:
                data_lines = [line.strip() for line in data_match.group(1).strip().split('\n') if line.strip()]
                for row_idx, line in enumerate(data_lines, 2):
                    data_values = [val.strip() for val in line.split(',')]
                    for col_idx, value in enumerate(data_values, 1):
                        if col_idx <= len(columns):
                            ws.cell(row=row_idx, column=col_idx, value=value)
            
            # Save file
            timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
            safe_name = re.sub(r"[^A-Za-z0-9_\- ]", "_", sheet_name).strip().replace(" ", "_")
            filename = f"{safe_name}-{timestamp}.xlsx"
            file_path = self.sheets_dir / filename
            wb.save(str(file_path))
            
            public_url = f"{self.public_base_url}/files/{filename}"
            
            return {
                "success": True,
                "message": f"✅ Excel file created: {sheet_name}",
                "file_path": str(file_path),
                "public_url": public_url,
                "filename": filename,
                "sheet_name": sheet_name,
                "columns": columns
            }
            
        except Exception as e:
            logger.error(f"Error creating Excel file: {e}")
            return {
                "success": False,
                "message": f"❌ Error creating Excel file: {str(e)}"
            }
    
    def list_files(self, file_type: str = None) -> List[Dict[str, Any]]:
        """List files in the cloud output directory."""
        files = []
        for file_path in self.cloud_out_dir.iterdir():
            if file_path.is_file():
                if file_type is None or file_path.suffix.lower() == f'.{file_type.lower()}':
                    files.append({
                        "name": file_path.name,
                        "path": str(file_path),
                        "size": file_path.stat().st_size,
                        "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                        "public_url": f"{self.public_base_url}/files/{file_path.name}"
                    })
        return sorted(files, key=lambda x: x["modified"], reverse=True)
    
    def get_file_info(self, filename: str) -> Dict[str, Any]:
        """Get information about a specific file."""
        file_path = self.cloud_out_dir / filename
        if not file_path.exists():
            return {"error": "File not found"}
        
        return {
            "name": file_path.name,
            "path": str(file_path),
            "size": file_path.stat().st_size,
            "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
            "public_url": f"{self.public_base_url}/files/{file_path.name}",
            "exists": True
        }

# Global file manager instance
file_manager = FileManager()

async def demo_file_operations():
    """Demonstrate file operations."""
    print("📁 File Management Demo")
    print("=" * 50)
    
    # Test 1: Create presentation
    print("\n1. Creating presentation...")
    result = await file_manager.process_file_request("create ppt about AI trends 2024")
    print(f"Result: {result['message']}")
    if result['success']:
        print(f"File: {result['filename']}")
        print(f"URL: {result['public_url']}")
    
    # Test 2: List files
    print("\n2. Listing files...")
    files = file_manager.list_files('pptx')
    for file in files[:3]:  # Show first 3
        print(f"  - {file['name']} ({file['size']} bytes)")
    
    # Test 3: Create Excel
    print("\n3. Creating Excel file...")
    result = await file_manager.process_file_request("create excel with sales data for Q1")
    print(f"Result: {result['message']}")

if __name__ == "__main__":
    asyncio.run(demo_file_operations())
